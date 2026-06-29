"""
Extractor universal de documentos: PDF, Word, Excel, PowerPoint.

Detecta el tipo de archivo por extensión o MIME type y extrae texto plano
o markdown para integración en AGENTS.md, README.md y context.md.
"""

import io
import re

from browser_api import hacer_get


def extraer_texto_documento(url: str) -> str | None:
    """
    Extrae texto de cualquier documento soportado (PDF, DOCX, XLSX, PPTX).

    Args:
        url: URL del documento en Moodle (pluginfile.php)

    Returns:
        Texto extraído formateado, o None si falló / no soportado.
    """
    try:
        contenido = hacer_get(url)
        stream = io.BytesIO(contenido)

        # Detectar tipo por extensión en URL
        url_lower = url.lower()
        if ".pdf" in url_lower:
            return _extraer_pdf(stream)
        elif ".docx" in url_lower or ".doc" in url_lower:
            return _extraer_word(stream)
        elif ".xlsx" in url_lower or ".xls" in url_lower:
            return _extraer_excel(stream)
        elif ".pptx" in url_lower or ".ppt" in url_lower:
            return _extraer_ppt(stream)
        else:
            # Fallback: intentar como PDF (algunos URLs no tienen extensión)
            return _extraer_pdf(stream)

    except Exception as e:
        print(f"[WARN] Error extrayendo {url}: {e}")
        return None


def _limpiar_texto(texto: str) -> str:
    """Normaliza espacios y líneas."""
    texto = re.sub(r'\n{3,}', '\n\n', texto)
    texto = re.sub(r' {2,}', ' ', texto)
    return texto.strip()


def _extraer_pdf(stream: io.BytesIO) -> str | None:
    import fitz
    doc = fitz.open(stream=stream, filetype="pdf")
    partes = []
    for pagina in doc:
        txt = pagina.get_text()
        if txt.strip():
            partes.append(txt)
    doc.close()
    return _limpiar_texto("\n\n".join(partes)) if partes else None


def _extraer_word(stream: io.BytesIO) -> str | None:
    from docx import Document
    doc = Document(stream)
    partes = []
    for para in doc.paragraphs:
        if para.text.strip():
            partes.append(para.text)
    return _limpiar_texto("\n\n".join(partes)) if partes else None


def _extraer_excel(stream: io.BytesIO) -> str | None:
    from openpyxl import load_workbook
    wb = load_workbook(stream, data_only=True, read_only=True)
    lineas = []
    for sheet_name in wb.sheetnames:
        lineas.append(f"## Hoja: {sheet_name}")
        ws = wb[sheet_name]
        for row in ws.iter_rows(values_only=True):
            # Filtrar celdas vacías
            celdas = [str(c) if c is not None else "" for c in row]
            celdas = [c for c in celdas if c.strip()]
            if celdas:
                lineas.append(" | ".join(celdas))
        lineas.append("")
    wb.close()
    return _limpiar_texto("\n".join(lineas)) if lineas else None


def _extraer_ppt(stream: io.BytesIO) -> str | None:
    from pptx import Presentation
    prs = Presentation(stream)
    diapositivas = []
    for i, slide in enumerate(prs.slides, 1):
        textos = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                textos.append(shape.text.strip())
        if textos:
            diapositivas.append(f"### Diapositiva {i}\n\n" + "\n\n".join(textos))
    return _limpiar_texto("\n\n".join(diapositivas)) if diapositivas else None
